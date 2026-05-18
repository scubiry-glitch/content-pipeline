#!/usr/bin/env python3
"""Generate PPT for 租赁好房子·评审通过后汇报 — v3 clean design"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ── minimal, professional
NAVY      = RGBColor(0x0F, 0x2B, 0x5B)   # 深藏蓝 (primary)
BLUE      = RGBColor(0x1A, 0x6B, 0xB5)   # 蓝 (accent)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF9)   # 页面底色
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1A, 0x1A, 0x2E)
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x5A)
MID_GRAY  = RGBColor(0x8A, 0x8A, 0x9A)
LIGHT_GRAY= RGBColor(0xE8, 0xEB, 0xF0)
GREEN     = RGBColor(0x0E, 0x8A, 0x5E)   # positive/done
RED       = RGBColor(0xD0, 0x3B, 0x3B)   # alert/risk
GOLD      = RGBColor(0xB8, 0x86, 0x1A)   # caution
PURPLE    = RGBColor(0x5C, 0x4E, 0xD0)   # direction2 accent

# ── Paths ──
SCRIPT_DIR  = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR  = os.path.join(SCRIPT_DIR, 'ppt-assets')
OUTPUT_PATH = os.path.join(SCRIPT_DIR, '租赁好房子·评审通过后汇报.pptx')

# ── Presentation ──
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
# ── Helper functions ──

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def tx(slide, left, top, w, h, text, sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    """Add a text box — the workhorse"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return box

def rect(slide, left, top, w, h, fill, border=None, border_w=1):
    """Add a rectangle shape"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(border_w)
    else: s.line.fill.background()
    return s

def accent_line(slide, left, top, w, color=BLUE, thickness=0.04):
    """Thin colored accent bar"""
    return rect(slide, left, top, w, thickness, color)

def section_title(slide, number, title, subtitle=None):
    """Standard page header: number badge + title + optional subtitle"""
    # Accent line at very top
    rect(slide, 0, 0, 13.333, 0.06, NAVY)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(0.45), Inches(0.55), Inches(0.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = NAVY; circle.line.fill.background()
    tf = circle.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(number); p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Inches(0.05)
    # Title
    tx(slide, 1.45, 0.4, 8, 0.6, title, sz=28, bold=True, color=NAVY)
    if subtitle:
        tx(slide, 1.45, 0.95, 10, 0.35, subtitle, sz=14, color=MID_GRAY)

def card(slide, left, top, w, h, title, lines, accent=BLUE, fill=WHITE):
    """Clean card with left accent bar"""
    # Card body
    s = rect(slide, left, top, w, h, fill)
    # Left accent bar
    rect(slide, left, top, 0.06, h, accent)
    # Text
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAVY; p.font.name = 'Microsoft YaHei'
    p.space_after = Pt(6)
    for line in lines:
        p2 = tf.add_paragraph(); p2.text = line
        p2.font.size = Pt(11); p2.font.color.rgb = DARK_GRAY; p2.font.name = 'Microsoft YaHei'
        p2.space_after = Pt(2)
    return s

def img(slide, path, left, top, w=None, h=None):
    """Add image with fallback placeholder"""
    full = os.path.join(ASSETS_DIR, path) if not os.path.isabs(path) else path
    try:
        if os.path.exists(full):
            kw = {}
            if w: kw['width'] = Inches(w)
            if h: kw['height'] = Inches(h)
            slide.shapes.add_picture(full, Inches(left), Inches(top), **kw)
            return
    except Exception:
        pass
    s = rect(slide, left, top, w or 4, h or 2.5, LIGHT_GRAY, LIGHT_GRAY)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"[截图]\n{path}"; p.font.size = Pt(11); p.font.color.rgb = MID_GRAY; p.alignment = PP_ALIGN.CENTER
# ═══════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

# Geometric accent shapes
rect(slide, 0, 0, 0.15, 7.5, BLUE)           # left stripe
rect(slide, 0, 3.8, 13.333, 0.02, BLUE)       # horizontal divider
rect(slide, 12.5, 0, 0.833, 3.8, RGBColor(0x14,0x3D,0x72))  # right accent

# Title block
tx(slide, 1.0, 1.8, 10, 1.2, '租赁好房子 · 评审通过后汇报', sz=44, bold=True, color=WHITE)
tx(slide, 1.0, 3.0, 10, 0.6, '从"纸面标准"到"业务武器"', sz=22, color=RGBColor(0x8A,0xBB,0xDE))

# Below divider
tx(slide, 1.0, 4.3, 10, 0.5, '2026年5月', sz=15, color=RGBColor(0x6A,0x9B,0xCD))
tx(slide, 1.0, 5.0, 10, 0.4, '标准已立 · 闭环才是护城河', sz=16, color=RGBColor(0x5A,0x8B,0xBD))
# ═══════════════════════════════════════════════
# SLIDE 2: WHY
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 1, '为什么做', '标准通过了，但标准本身不产生价值')

# Key statement
tx(slide, 0.7, 1.6, 11, 0.5, '对手抄走"舒适+科技+绿色"只需三天，抄走24万条数据+80万套评分体系却要三年', sz=16, bold=True, color=NAVY)
tx(slide, 0.7, 2.1, 11, 0.4, '谁完成闭环，谁就是行业基础设施；谁没完成，标准就是一张纸', sz=13, color=DARK_GRAY)

# Three columns — clean cards with accent bars
card(slide, 0.7, 2.8, 3.8, 3.8,
     '不做 → 被定义', [
         '标准没有护城河',
         '对手能抄走定义',
         '但抄不走数据+评分',
         '',
         '谁没完成闭环',
         '标准=一张纸',
     ], accent=GOLD)

card(slide, 4.8, 2.8, 3.8, 3.8,
     '做了 → 定义别人', [
         '方法论 → 定义和定价',
         '工具(AI) → 边际趋零',
         '认证 → 人才锁定',
         '',
         '从5%净利率中介',
         '升级为行业操作系统',
     ], accent=GREEN)

card(slide, 8.9, 2.8, 3.8, 3.8,
     '晚做 → 窗口关闭', [
         '标准话语权有时效',
         'APEC 12月须有案例',
         '人社部需1W人次',
         'AI数据壁垒在贬值',
         '',
         '不做=被定义',
         '做了=定义别人，晚做=关闭',
     ], accent=RED)
# ═══════════════════════════════════════════════
# SLIDE 3: IRON TRIANGLE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 2, '铁三角', '方法论 + 工具 + 职业生涯 = 行业操作系统')

# Three pillars
card(slide, 0.7, 1.5, 3.8, 3.2,
     '📐 方法论', [
         '好房子评价标准',
         'HQI 量化评分',
         'CAPM+Hedonic 定价模型',
         'GWR+LightGBM 空间量化',
     ], accent=NAVY)

card(slide, 4.8, 1.5, 3.8, 3.2,
     '🔧 工具', [
         'AI评级（图片→评分）',
         'REFA 一线作业',
         'PL&UE 看板',
         '会议纪要应用',
     ], accent=BLUE)

card(slide, 8.9, 1.5, 3.8, 3.2,
     '👤 职业生涯', [
         'REFP认证（对标CFP）',
         'LV1资管→LV2战队长→LV3总监',
         '人社部从业资格（1W人次后）',
         '清华五道口课程研发',
     ], accent=PURPLE)

# Bidirectional logic
tx(slide, 1.5, 4.8, 3.0, 0.3, '标准定义评什么 / 模型定义怎么评 →', sz=10, color=BLUE, align=PP_ALIGN.CENTER)
tx(slide, 5.5, 4.8, 2.5, 0.3, '← 工具提效 / 数据量化 →', sz=10, color=BLUE, align=PP_ALIGN.CENTER)
tx(slide, 9.5, 4.8, 2.8, 0.3, '← 认证落地到人', sz=10, color=BLUE, align=PP_ALIGN.CENTER)

# Bottom: 6个1 progress
accent_line(slide, 0.7, 5.3, 12.0, LIGHT_GRAY, 0.03)
tx(slide, 0.7, 5.5, 2, 0.3, '"6个1"进度', sz=13, bold=True, color=NAVY)

items = [
    ('1 行业标准', '中房协', '✅', GREEN),
    ('2 职业认证', '人社部', '🔄', BLUE),
    ('3 发布会', '人民日报', '📋', MID_GRAY),
    ('4 白皮书', '清华', '📋', MID_GRAY),
    ('5 标杆项目', '—', '🔄', BLUE),
    ('6 AI工具', '—', '🔄', BLUE),
]
for i, (name, partner, status, color) in enumerate(items):
    left = 0.7 + i * 2.05
    rect(slide, left, 5.9, 1.9, 1.0, WHITE, LIGHT_GRAY, 0.5)
    tx(slide, left+0.1, 5.95, 1.7, 0.3, name, sz=11, bold=True, color=NAVY)
    tx(slide, left+0.1, 6.25, 1.7, 0.25, f'{partner}  {status}', sz=10, color=color)

# Bottom support row
tx(slide, 0.7, 7.0, 12, 0.3, '支撑：📚 学术(清华论文·6月) · 🏛️ 机构(房协+人社部) · 📊 数据(80万套+1000套+24万条)', sz=11, color=MID_GRAY)
# ═══════════════════════════════════════════════
# SLIDE 4: DIRECTION 1 — 标准武器化
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 3, '方向一：标准武器化', '定义好 → 量化好 → 定价好 · 对手抄标准，抄不走闭环')

# Three-layer architecture — left
card(slide, 0.7, 1.5, 3.5, 1.4,
     '好房子评价标准', ['定义"什么是好" · 协会背书 · 原则性'], accent=NAVY)
# Arrow
tx(slide, 2.0, 2.95, 1.5, 0.3, '↓  原则→指标', sz=10, color=BLUE, align=PP_ALIGN.CENTER)
card(slide, 0.7, 3.2, 3.5, 1.4,
     'HQI', ['量化"有多好" · 可打分·可验证'], accent=BLUE)
tx(slide, 2.0, 4.65, 1.5, 0.3, '↓  评分→溢价', sz=10, color=BLUE, align=PP_ALIGN.CENTER)
card(slide, 0.7, 4.9, 3.5, 1.4,
     'CAPM+Hedonic', ['计算"值多少钱" · 可定价·可预测 · 24万条'], accent=BLUE)

# AI bridge — center
card(slide, 4.6, 2.2, 2.8, 3.6,
     '🤖 AI评级', [
         '规模化桥梁',
         '',
         '拍照→AI评分→报告',
         '成本趋零·全覆盖',
         '',
         '标准定义评什么',
         'HQI提供评分框架',
         'AI实现规模化执行',
     ], accent=BLUE, fill=RGBColor(0xEE,0xF3,0xFA))

# Right: pilot + moat
card(slide, 7.8, 1.5, 2.5, 2.0,
     '🎯 美租试点', [
         '1000+套·最掐尖',
         '3个月出模型',
         '三星+15%·五星+30%',
     ], accent=GREEN)

card(slide, 7.8, 3.8, 2.5, 2.0,
     '🔒 护城河', [
         '抄不走HQI评分',
         '抄不走24万条数据',
         '抄不走AI模型',
     ], accent=GREEN)

# Screenshots — right side, larger
tx(slide, 10.5, 1.4, 2.7, 0.3, 'AI评级演示', sz=12, bold=True, color=NAVY)
img(slide, 'upload-room-photos.png', 10.5, 1.8, w=2.6)
img(slide, 'detailed-quotation.png', 10.5, 4.4, w=2.6)

# Bottom: CAPM table + formula
accent_line(slide, 0.7, 6.5, 12.0, LIGHT_GRAY, 0.02)

# CAPM condensed
capm = [('Hedonic', '变量边际效应 → 定价底层逻辑'), ('GWR', '空间异质性 → 差异化策略'), ('LightGBM', '高精度预测 → 自动定价引擎'), ('CAPM', 'ROI = [CapM(after)-CapM(before)] / (托管费+装配成本)')]
for i, (name, desc) in enumerate(capm):
    left = 0.7 + i * 3.1
    tx(slide, left, 6.6, 2.9, 0.25, name, sz=11, bold=True, color=NAVY)
    tx(slide, left, 6.85, 2.9, 0.25, desc, sz=10, color=DARK_GRAY)

# Risk callout
rect(slide, 0.7, 7.15, 12.0, 0.25, RGBColor(0xFD,0xF0,0xF0))
tx(slide, 0.9, 7.15, 11.5, 0.25, '⚠️ 必须用外部验证校准AI评分，不能用内部数据自证', sz=10, bold=True, color=RED)
# ═══════════════════════════════════════════════
# SLIDE 5: DIRECTION 2 — 工具落地闭环
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 4, '方向二：工具落地闭环', '工具→人→项目→新场景 · 四环咬合')

# Ring 1: REFA/NPG
card(slide, 0.7, 1.5, 3.6, 2.6,
     '第一环：REFA·NPG', [
         'N 一线 — 信号采集器',
         '  录音记录真实场景',
         'P 总监 — 实验室',
         '  基于AI输出做复盘',
         'G 总部 — 研发中心',
         '  聚合数据→领域包',
     ], accent=BLUE)
tx(slide, 0.7, 4.2, 3.6, 0.3, '工具让人有力', sz=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Arrow
tx(slide, 4.35, 2.5, 0.4, 0.3, '→', sz=20, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Ring 2: People
card(slide, 4.8, 1.5, 3.6, 2.6,
     '第二环：先行人群', [
         '资管经理 → AI工具先行',
         '  REFA解决信息不对称',
         '设计师 → 认证先行',
         '  "认证设计师"工牌',
         '两者交汇 → 人+标准+工具',
     ], accent=GREEN)
tx(slide, 4.8, 4.2, 3.6, 0.3, '人让标准落地', sz=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Arrow
tx(slide, 8.45, 2.5, 0.4, 0.3, '→', sz=20, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Ring 3+4: Projects + High-end
card(slide, 8.9, 1.5, 3.8, 1.2,
     '第三环：样板项目', [
         '杭州→建造验证 · 南京→运营验证 · 上海FAM→金融闭环',
     ], accent=PURPLE)

card(slide, 8.9, 2.85, 3.8, 1.2,
     '第四环：高端租赁', [
         '美租第二品牌 · 以客找房 · B2B获客',
     ], accent=GOLD)

# Screenshots row
accent_line(slide, 0.7, 4.6, 12.0, LIGHT_GRAY, 0.02)
tx(slide, 0.7, 4.7, 3, 0.3, 'NPG系统截图', sz=12, bold=True, color=NAVY)

img(slide, 'npg-matrix-overview.png', 0.7, 5.1, w=3.5)
img(slide, 'npg-recording.png', 4.5, 5.1, w=3.5)
img(slide, 'npg-experiment.png', 8.3, 5.1, w=3.5)

# Risk bar
rect(slide, 0.7, 7.15, 12.0, 0.25, RGBColor(0xFD,0xF0,0xF0))
tx(slide, 0.9, 7.15, 11.5, 0.25, '⚠️ 必须给具体场景，不能做"自嗨工程"', sz=10, bold=True, color=RED)
# ═══════════════════════════════════════════════
# SLIDE 6: DIRECTION 3 — 公众防线
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 5, '方向三：公众防线', '面向公众的独立议题 · 先守住底线，再谈上限')

# WHY
rect(slide, 0.7, 1.5, 12.0, 0.7, RGBColor(0xFD,0xF4,0xE5))
tx(slide, 0.9, 1.55, 11.5, 0.6, '80万套规模=80万倍杠杆 · 一个燃气事故的舆论风险可以抵消所有品牌建设', sz=14, bold=True, color=GOLD)

# Bottom line: Responsibility
card(slide, 0.7, 2.5, 5.7, 3.5,
     '🛡️ 底线：责任界定', [
         '当前"大包大揽"，责任边界模糊=无限风险',
         '',
         '推动"省心租管理责任清单"：',
         '',
         '贝壳 → 招租时效、租金代扣、维修响应',
         '物业/燃气 → 燃气管道、电梯维保、公共安全',
         '用户 → 使用不当、人为损坏、违规改造',
         '',
         '⏰ 1个月出初稿 · 法务+运营先上',
     ], accent=RED)

# Upper limit: Value
card(slide, 6.7, 2.5, 5.9, 3.5,
     '🌟 上限：价值传达', [
         '"舒适+科技+绿色"要用租客语言：',
         '',
         '专家语言                    租客语言',
         '好房子标准达标      →    "不满意7天可退"',
         'AI智能匹配            →    "3分钟匹配最适合"',
         'ENF级环保材料      →    "检测报告扫码可查"',
         '',
         '需品宣+业务对齐 · 建议A/B测试',
     ], accent=GREEN)

# Priority bar
rect(slide, 0.7, 6.2, 12.0, 0.5, RGBColor(0xEE,0xF3,0xFA))
tx(slide, 0.9, 6.25, 11.5, 0.4, '优先级：责任界定 > 价值传达 · 法务+运营先上，品宣后上', sz=13, bold=True, color=NAVY)

# Direction arrow
tx(slide, 3.0, 5.95, 5, 0.3, '守住底线 ────────→ 才敢谈上限', sz=12, color=BLUE, align=PP_ALIGN.CENTER)
# ═══════════════════════════════════════════════
# SLIDE 7: TIMELINE + DECISIONS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 6, '时间线与决策', '4项必须拍板 · 2项确认')

# Timeline — top half, simplified horizontal bars
milestones = [
    ('5月', 'REFP大纲完成', GREEN),
    ('6月', '论文+机构确认', BLUE),
    ('7月', '100人次考试', BLUE),
    ('11月', '认证发布会', MID_GRAY),
    ('12月', 'APEC论坛', MID_GRAY),
    ('次年', '1W人次→人社部', MID_GRAY),
]
# Single timeline bar
rect(slide, 0.7, 1.65, 12.0, 0.08, LIGHT_GRAY)
for i, (time, desc, color) in enumerate(milestones):
    left = 0.7 + i * 2.05
    # Dot on bar
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left+0.85), Inches(1.57), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    # Text below
    tx(slide, left, 1.9, 1.9, 0.25, time, sz=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    tx(slide, left, 2.15, 1.9, 0.4, desc, sz=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

accent_line(slide, 0.7, 2.7, 12.0, LIGHT_GRAY, 0.02)

# Decisions — bottom half, 2x2 grid
decisions = [
    ('1', '批准2-3个集中式样板项目立项', '6个月完成装修+入住+AI评级 → APEC展示案例', NAVY),
    ('2', '批准AI评级在美租试点', '1000+套训练集 · 3个月出模型', BLUE),
    ('3', '批准省心租责任清单制定', '法务+运营+品宣 · 1个月初稿', RED),
    ('4', '确认铁三角推进路径', '方法论+工具+职业生涯 · 7月→11月→次年', PURPLE),
]
for i, (num, title, detail, color) in enumerate(decisions):
    col = i % 2; row = i // 2
    left = 0.7 + col * 6.2; top = 2.9 + row * 2.0
    # Card
    rect(slide, left, top, 5.9, 1.7, WHITE, LIGHT_GRAY, 0.5)
    # Left accent bar
    rect(slide, left, top, 0.06, 1.7, color)
    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left+0.2), Inches(top+0.15), Inches(0.4), Inches(0.4))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    tf_c = circle.text_frame; tf_c.word_wrap = False; tf_c.margin_top = Inches(0.02)
    p = tf_c.paragraphs[0]; p.text = num; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    # Title + detail
    tx(slide, left+0.75, top+0.15, 4.8, 0.4, title, sz=14, bold=True, color=NAVY)
    tx(slide, left+0.75, top+0.6, 4.8, 0.5, detail, sz=11, color=DARK_GRAY)

# Bottom: 2 items to confirm
rect(slide, 0.7, 7.05, 12.0, 0.35, LIGHT_GRAY)
tx(slide, 0.9, 7.07, 11.5, 0.3, '📋 需确认但不必当天：⑤ 三层架构定位  ⑥ 价值传达转租客语言', sz=11, color=MID_GRAY)
# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
prs.save(OUTPUT_PATH)
print(f'✅ Saved: {OUTPUT_PATH}')
print(f'   Slides: {len(prs.slides)}')
